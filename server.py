# server.py
import os
import base64
import json
import requests
import re
import io
import time
import pytz
import asyncio
import edge_tts
import traceback
import random
from github import Github
from google import genai
from google.genai import types
from datetime import datetime
from fastapi import FastAPI, Header, HTTPException, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional, Union

# --- LIBRARIES ---
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Pt as PptxPt, Inches
from pptx.dml.color import RGBColor as PptxColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from duckduckgo_search import DDGS

# --- KEYS ---
SERVER_SECRET_KEY = "Mbuso.08@"
ELEVENLABS_API_KEY = os.environ.get("ELEVENLABS_API_KEY")
BRIAN_VOICE_ID = "nPczCjzI2devNBz1zQrb"
HUGGINGFACE_API_KEY = os.environ.get("HUGGINGFACE_API_KEY")
GITHUB_TOKEN = os.environ.get("GITHUB_TOKEN")
GITHUB_REPO_NAME = os.environ.get("GITHUB_REPO")
MEMORY_FILE = "alfred_memory.json"

# --- CLOUD MEMORY ---
def pull_memory_from_cloud():
    if not GITHUB_TOKEN or not GITHUB_REPO_NAME: return
    try:
        g = Github(GITHUB_TOKEN)
        repo = g.get_repo(GITHUB_REPO_NAME)
        contents = repo.get_contents(MEMORY_FILE)
        remote_data = json.loads(contents.decoded_content.decode())
        with open(MEMORY_FILE, "w") as f: json.dump(remote_data, f, indent=4)
    except: pass

def push_memory_to_cloud():
    if not GITHUB_TOKEN or not GITHUB_REPO_NAME: return
    try:
        g = Github(GITHUB_TOKEN)
        repo = g.get_repo(GITHUB_REPO_NAME)
        contents = repo.get_contents(MEMORY_FILE)
        with open(MEMORY_FILE, "r") as f: local_data = json.load(f)
        repo.update_file(contents.path, "Alfred Memory Update", json.dumps(local_data, indent=4), contents.sha)
    except: pass

def load_memory():
    if not os.path.exists(MEMORY_FILE): return "No prior knowledge."
    try:
        with open(MEMORY_FILE, "r") as f:
            data = json.load(f)
            return "\n".join([f"- {fact}" for fact in data.get("facts", [])])
    except: return "Memory Corrupted."

def save_memory_fact(fact):
    try:
        data = {"facts": []}
        if os.path.exists(MEMORY_FILE):
            with open(MEMORY_FILE, "r") as f: data = json.load(f)
        if fact not in data["facts"]:
            data["facts"].append(fact)
            with open(MEMORY_FILE, "w") as f: json.dump(data, f, indent=4)
            try:
                loop = asyncio.get_event_loop()
                loop.run_in_executor(None, push_memory_to_cloud)
            except: pass
    except: pass

pull_memory_from_cloud()

# --- SYSTEM PROMPT ---
def get_system_instructions():
    return f"""
You are Alfred, an elite intelligent assistant.
User Name: Mbuso.
ADDRESS USER AS: "Sir".

*** MEMORY ***
{load_memory()}

*** IOS COMMANDS (STRICT JSON) ***
1. REMINDER: <<<IOS_CMD>>> {{"type": "reminder", "title": "Title", "time": "HH:MM"}} <<<END>>>
2. ALARM: <<<IOS_CMD>>> {{"type": "alarm", "time": "HH:MM", "label": "Name"}} <<<END>>>
3. MESSAGE: <<<IOS_CMD>>> {{"type": "message", "contact": "Name", "body": "Text"}} <<<END>>>
4. HEALTH: <<<IOS_CMD>>> {{"type": "health", "metric": "steps", "value": "100"}} <<<END>>>

*** PROTOCOLS ***
- FILES: <<<FILE_START_PPTX>>> ... <<<FILE_END>>>
- IMAGES: "IMAGE_GEN_REQUEST: [Prompt]"
- DEFAULT: Concise answer. Save facts: <<<MEM_SAVE>>> fact <<<MEM_END>>>.
"""

api_key = os.environ.get("GEMINI_API_KEY")
client = None
if api_key: client = genai.Client(api_key=api_key)

app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=False, allow_methods=["*"], allow_headers=["*"])

class ChatMessage(BaseModel):
    role: str
    content: str

class UploadedFile(BaseModel):
    data: str
    mime_type: str

class UserContext(BaseModel):
    local_time: Optional[str] = None
    location: Optional[str] = None
    health_stats: Optional[str] = None

class UserRequest(BaseModel):
    command: str
    files: List[UploadedFile] = []
    context: Optional[Union[UserContext, str]] = None
    speak: bool = False
    thinking_mode: bool = False
    history: List[ChatMessage] = []

# --- TOOLS ---
def perform_web_search(query):
    try:
        results = DDGS().text(query, max_results=3)
        if not results: return None
        summary = "REAL-TIME SEARCH DATA:\n"
        for r in results: summary += f"- {r['title']}: {r['body']}\n"
        return summary
    except: return None

def get_url_content(text):
    url_match = re.search(r'(https?://[^\s]+)', text)
    if not url_match: return None
    url = url_match.group(0)
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        r = requests.get(url, headers=headers, timeout=10)
        s = BeautifulSoup(r.content, 'html.parser')
        for x in s(["script", "style"]): x.extract()
        return f"[WEB]: {s.get_text(separator=' ', strip=True)[:8000]}"
    except: return None

async def generate_voice_dual(text):
    clean = re.sub(r'<<<.*?>>>', '', text).replace('IMAGE_GEN_REQUEST:', '')
    if not clean.strip(): return None
    if ELEVENLABS_API_KEY:
        try:
            url = f"https://api.elevenlabs.io/v1/text-to-speech/{BRIAN_VOICE_ID}"
            headers = {"xi-api-key": ELEVENLABS_API_KEY, "Content-Type": "application/json"}
            res = requests.post(url, json={"text": clean[:1000], "model_id": "eleven_monolingual_v1"}, headers=headers)
            if res.status_code == 200: return base64.b64encode(res.content).decode('utf-8')
        except: pass
    try:
        communicate = edge_tts.Communicate(clean, "en-GB-RyanNeural")
        fname = f"voice_{int(time.time())}.mp3"
        await communicate.save(fname)
        with open(fname, "rb") as f: audio = f.read()
        os.remove(fname)
        return base64.b64encode(audio).decode('utf-8')
    except: return None

# --- VISUAL ENGINE ---
def generate_image(prompt):
    seed = int(time.time())
    encoded = requests.utils.quote(prompt)
    models = ["flux", "turbo"]
    for m in models:
        try:
            url = f"https://image.pollinations.ai/prompt/{encoded}?seed={seed}&width=1280&height=720&model={m}"
            r = requests.get(url, timeout=20)
            if r.status_code == 200: return base64.b64encode(r.content).decode('utf-8'), f"gen_{seed}.jpg"
        except: continue
    if HUGGINGFACE_API_KEY:
        try:
            h = {"Authorization": f"Bearer {HUGGINGFACE_API_KEY}"}
            r = requests.post("https://api-inference.huggingface.co/models/stabilityai/stable-diffusion-xl-base-1.0", headers=h, json={"inputs": prompt}, timeout=20)
            if r.status_code == 200: return base64.b64encode(r.content).decode('utf-8'), "gen.jpg"
        except: pass
    return None, "Busy"

# --- THEME ENGINE ---
def apply_neon_theme(slide, is_title=False):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = PptxColor(5, 5, 16)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top.fill.solid(); top.fill.fore_color.rgb = PptxColor(0, 243, 255); top.line.fill.background()
    btm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(10), Inches(0.15))
    btm.fill.solid(); btm.fill.fore_color.rgb = PptxColor(188, 19, 254); btm.line.fill.background()
    if is_title: card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3.5))
    else: card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    card.fill.solid(); card.fill.fore_color.rgb = PptxColor(30, 30, 46); card.line.color.rgb = PptxColor(60, 60, 80)

def create_file(content, ftype):
    try:
        buf = io.BytesIO()
        fname = f"Alfred_Export_{int(time.time())}"
        if ftype == "pptx":
            prs = Presentation()
            chunks = re.split(r'Slide \d+:', content, flags=re.IGNORECASE)
            s1 = prs.slides.add_slide(prs.slide_layouts[6]); apply_neon_theme(s1, True)
            title = s1.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(7.6), Inches(2)); title.text_frame.text = chunks[0].replace("Title:", "").strip()[:80]
            title.text_frame.paragraphs[0].font.size = PptxPt(44); title.text_frame.paragraphs[0].font.bold = True; title.text_frame.paragraphs[0].font.color.rgb = PptxColor(255, 255, 255); title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for c in chunks[1:]:
                lines = c.strip().split('\n')
                s = prs.slides.add_slide(prs.slide_layouts[6]); apply_neon_theme(s, False)
                h = s.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8)); h.text_frame.text = lines[0].strip()
                h.text_frame.paragraphs[0].font.size = PptxPt(32); h.text_frame.paragraphs[0].font.bold = True; h.text_frame.paragraphs[0].font.color.rgb = PptxColor(255, 255, 255)
                b = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5)); b.text_frame.word_wrap = True
                for line in lines[1:]:
                    if not line.strip(): continue
                    p = b.text_frame.add_paragraph(); p.text = "• " + line.lstrip("-*• ").strip()
                    p.font.size = PptxPt(20); p.font.color.rgb = PptxColor(220, 220, 220); p.space_after = PptxPt(14)
            prs.save(buf); return base64.b64encode(buf.getvalue()).decode('utf-8'), f"{fname}.pptx"
        elif ftype == "docx":
            d = Document(); d.add_paragraph(content); d.save(buf)
            return base64.b64encode(buf.getvalue()).decode('utf-8'), f"{fname}.docx"
        elif ftype == "txt":
            return base64.b64encode(content.encode('utf-8')).decode('utf-8'), f"{fname}.txt"
    except: return None, None

# --- ROUTE ---
@app.post("/command")
async def process_command(request: UserRequest, x_alfred_auth: Optional[str] = Header(None), x_client_device: Optional[str] = Header(None)):
    if x_alfred_auth != SERVER_SECRET_KEY: raise HTTPException(status_code=401)
    if not client: return {"response": "No API Key."}

    parsed_context = None
    if request.context:
        if isinstance(request.context, str):
            try:
                data = json.loads(request.context)
                parsed_context = UserContext(**data)
            except: parsed_context = None
        else:
            parsed_context = request.context

    sa_time = datetime.now(pytz.timezone('Africa/Johannesburg')).strftime("%H:%M")
    phone_context = f"\n[PHONE DATA -> Time: {parsed_context.local_time}, Loc: {parsed_context.location}]" if parsed_context else ""
    
    ctx = ""
    if any(x in request.command.lower() for x in ["search", "find", "news"]):
        res = perform_web_search(request.command)
        if res: ctx += f"\n{res}\n"
    
    prompt = f"[Current SAST: {sa_time}] {phone_context} {ctx} \nUser: {request.command}"
    parts = [prompt]
    if request.files:
        for f in request.files:
            try: parts.append(types.Part.from_bytes(data=base64.b64decode(f.data), mime_type=f.mime_type))
            except: pass

    # --- IOS MODE (NO STREAMING) ---
    if x_client_device == "ios":
        try:
            chat = client.chats.create(model='gemini-2.0-flash-thinking-exp-01-21', history=[types.Content(role="model" if m.role=="alfred" else "user", parts=[types.Part.from_text(text=m.content)]) for m in request.history], config=types.GenerateContentConfig(system_instruction=get_system_instructions(), tools=[types.Tool(google_search=types.GoogleSearch())]))
            response = chat.send_message(parts)
            text_reply = response.text
        except Exception as e:
            try:
                # SILENT FALLBACK (NO TEXT NOTIFICATION)
                chat = client.chats.create(model='gemini-2.0-flash', history=[types.Content(role="model" if m.role=="alfred" else "user", parts=[types.Part.from_text(text=m.content)]) for m in request.history], config=types.GenerateContentConfig(system_instruction=get_system_instructions(), tools=[types.Tool(google_search=types.GoogleSearch())]))
                response = chat.send_message(parts)
                text_reply = response.text
            except Exception as e2:
                return Response(content=f"Alfred Offline. Error: {str(e2)}", media_type="text/plain")

        if "<<<MEM_SAVE>>>" in text_reply:
            try: save_memory_fact(text_reply.split("<<<MEM_SAVE>>>")[1].split("<<<MEM_END>>>")[0].strip())
            except: pass
        
        return Response(content=text_reply, media_type="text/plain")

    # --- WEB MODE (STREAMING) ---
    async def stream_gen():
        try:
            chat = client.chats.create(model='gemini-2.0-flash-thinking-exp-01-21', history=[types.Content(role="model" if m.role=="alfred" else "user", parts=[types.Part.from_text(text=m.content)]) for m in request.history], config=types.GenerateContentConfig(system_instruction=get_system_instructions(), tools=[types.Tool(google_search=types.GoogleSearch())]))
            full_text = ""
            for chunk in chat.send_message_stream(parts):
                if chunk.text:
                    full_text += chunk.text
                    yield json.dumps({"type": "text", "content": chunk.text}) + "\n"
                    await asyncio.sleep(0.01)
        except Exception as e:
            if "429" in str(e) or "404" in str(e):
                # SILENT FALLBACK FOR WEB TOO (No "Switching" Text)
                chat = client.chats.create(model='gemini-2.0-flash', history=[types.Content(role="model" if m.role=="alfred" else "user", parts=[types.Part.from_text(text=m.content)]) for m in request.history], config=types.GenerateContentConfig(system_instruction=get_system_instructions(), tools=[types.Tool(google_search=types.GoogleSearch())]))
                for chunk in chat.send_message_stream(parts):
                    if chunk.text:
                        full_text += chunk.text
                        yield json.dumps({"type": "text", "content": chunk.text}) + "\n"
                        await asyncio.sleep(0.01)
            else:
                yield json.dumps({"type": "text", "content": str(e)}) + "\n"
        
        if "<<<MEM_SAVE>>>" in full_text:
            try: save_memory_fact(full_text.split("<<<MEM_SAVE>>>")[1].split("<<<MEM_END>>>")[0].strip())
            except: pass
            
        # Check for Images/Files/Voice
        final_payload = {"type": "meta", "audio": None, "file": None, "chart": None}
        
        img_match = re.search(r'IMAGE_GEN_REQUEST:\s*(.*)', full_text, re.IGNORECASE)
        if img_match:
            gf, gn = generate_image(img_match.group(1).strip())
            if gf: final_payload["file"] = {"data": gf, "name": gn}

        content, ftype = None, None
        if "<<<FILE_START_PPTX>>>" in full_text: content = full_text.split("<<<FILE_START_PPTX>>>")[1].split("<<<FILE_END>>>")[0].strip(); ftype="pptx"
        elif "<<<FILE_START_DOCX>>>" in full_text: content = full_text.split("<<<FILE_START_DOCX>>>")[1].split("<<<FILE_END>>>")[0].strip(); ftype="docx"
        
        if content:
            gf, gn = create_file(content, ftype)
            if gf: final_payload["file"] = {"data": gf, "name": gn}

        if "<<<CHART_START>>>" in full_text:
            try:
                js = full_text.split("<<<CHART_START>>>")[1].split("<<<CHART_END>>>")[0].strip().replace("```json", "").replace("```", "")
                final_payload["chart"] = json.loads(js)
            except: pass

        if request.speak:
            final_payload["audio"] = await generate_voice_dual(full_text)

        yield json.dumps(final_payload) + "\n"

    return StreamingResponse(stream_gen(), media_type="application/x-ndjson")
